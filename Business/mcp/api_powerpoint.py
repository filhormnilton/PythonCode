"""
MCP: api_powerpoint
Tools for creating, reading, updating, and deleting PowerPoint presentations.
"""
from pathlib import Path

from langchain_core.tools import tool

from Business.config import CONFIG


def _ensure_dir(directory: Path) -> None:
    directory.mkdir(parents=True, exist_ok=True)


def _pptx():
    try:
        from pptx import Presentation  # type: ignore
        from pptx.util import Inches, Pt  # type: ignore
        return Presentation, Inches, Pt
    except ImportError:
        return None, None, None


@tool
def create_presentation(filename: str, title: str, slides_content: str) -> str:
    """Create a new PowerPoint presentation.

    Args:
        filename: Name of the file (without extension).
        title: Title slide text.
        slides_content: Slide content separated by '---SLIDE---'. Each block may
                        contain a title line prefixed with 'TITLE:' and body text.

    Returns:
        Absolute path of the created .pptx file.
    """
    Presentation, Inches, Pt = _pptx()
    if Presentation is None:
        return "ERROR: python-pptx is not installed. Run: pip install python-pptx"

    _ensure_dir(CONFIG.output.slides_dir)
    path = CONFIG.output.slides_dir / f"{filename}.pptx"

    prs = Presentation()
    slide_layout_title = prs.slide_layouts[0]
    slide_layout_content = prs.slide_layouts[1]

    # Title slide
    slide = prs.slides.add_slide(slide_layout_title)
    slide.shapes.title.text = title
    if len(slide.placeholders) > 1:
        slide.placeholders[1].text = ""

    for block in slides_content.split("---SLIDE---"):
        block = block.strip()
        if not block:
            continue
        lines = block.splitlines()
        slide_title = ""
        body_lines = []
        for line in lines:
            if line.startswith("TITLE:"):
                slide_title = line[len("TITLE:"):].strip()
            else:
                body_lines.append(line)

        s = prs.slides.add_slide(slide_layout_content)
        s.shapes.title.text = slide_title
        tf = s.placeholders[1].text_frame
        tf.text = "\n".join(body_lines)

    prs.save(str(path))
    return str(path)


@tool
def read_presentation(filepath: str) -> str:
    """Read and return text content from a PowerPoint file.

    Args:
        filepath: Absolute or relative path to the .pptx file.

    Returns:
        Extracted text from all slides.
    """
    Presentation, _, _ = _pptx()
    if Presentation is None:
        return "ERROR: python-pptx is not installed. Run: pip install python-pptx"

    prs = Presentation(filepath)
    result = []
    for i, slide in enumerate(prs.slides, start=1):
        texts = [shape.text for shape in slide.shapes if shape.has_text_frame]
        result.append(f"[Slide {i}]\n" + "\n".join(texts))
    return "\n\n".join(result)


@tool
def add_slide(filepath: str, slide_title: str, slide_body: str) -> str:
    """Append a new slide to an existing presentation.

    Args:
        filepath: Absolute or relative path to the .pptx file.
        slide_title: Title for the new slide.
        slide_body: Body text for the new slide.

    Returns:
        Confirmation message with updated slide count.
    """
    Presentation, _, _ = _pptx()
    if Presentation is None:
        return "ERROR: python-pptx is not installed. Run: pip install python-pptx"

    prs = Presentation(filepath)
    slide_layout = prs.slide_layouts[1]
    s = prs.slides.add_slide(slide_layout)
    s.shapes.title.text = slide_title
    s.placeholders[1].text_frame.text = slide_body
    prs.save(filepath)
    return f"Slide added. Total slides: {len(prs.slides)}"


@tool
def update_slide(filepath: str, slide_index: int, slide_title: str, slide_body: str) -> str:
    """Update the title and body of an existing slide (0-based index).

    Args:
        filepath: Absolute or relative path to the .pptx file.
        slide_index: Zero-based index of the slide to update.
        slide_title: New title text.
        slide_body: New body text.

    Returns:
        Confirmation message.
    """
    Presentation, _, _ = _pptx()
    if Presentation is None:
        return "ERROR: python-pptx is not installed. Run: pip install python-pptx"

    prs = Presentation(filepath)
    if slide_index >= len(prs.slides):
        return f"ERROR: Slide index {slide_index} out of range (total: {len(prs.slides)})"
    slide = prs.slides[slide_index]
    if slide.shapes.title:
        slide.shapes.title.text = slide_title
    for shape in slide.shapes:
        if shape.has_text_frame and shape != slide.shapes.title:
            shape.text_frame.text = slide_body
            break
    prs.save(filepath)
    return f"Slide {slide_index} updated."


@tool
def delete_presentation(filepath: str) -> str:
    """Delete a PowerPoint file.

    Args:
        filepath: Absolute or relative path to the .pptx file.

    Returns:
        Confirmation message.
    """
    path = Path(filepath)
    if not path.exists():
        return f"ERROR: File not found: {filepath}"
    path.unlink()
    return f"Deleted: {filepath}"


POWERPOINT_TOOLS = [
    create_presentation,
    read_presentation,
    add_slide,
    update_slide,
    delete_presentation,
]
